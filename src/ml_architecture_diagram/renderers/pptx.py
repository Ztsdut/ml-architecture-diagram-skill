from __future__ import annotations

from pathlib import Path
import math

from ..layout import layout_figure
from ..theme import theme_for_spec
from ..publication_design import compile_publication_spec
from ..routing import fanin_bundle_geometry, fanout_bundle_geometry


def _hex_rgb(hexstr: str):
    from pptx.dml.color import RGBColor
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def render_pptx(spec: dict, output: str | Path) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        from pptx.oxml.xmlchemy import OxmlElement
    except ImportError as exc:
        raise RuntimeError("PPTX export requires python-pptx. Install with: pip install -e '.[export]'") from exc

    spec = compile_publication_spec(spec)
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    layout = layout_figure(spec)
    fig = spec.get("figure", {})
    theme = theme_for_spec(spec)
    font = fig.get("font", "Arial")
    node_map = {n["id"]: n for n in spec.get("nodes", [])}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    aspect = max(layout["height"] + 70, 520) / max(layout["width"], 800)
    prs.slide_height = Inches(max(7.5, min(13.333, 13.333 * aspect)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = Inches(0.25)
    top_title = Inches(0.42)
    usable_w = slide_w - 2 * margin
    usable_h = slide_h - Inches(0.72)
    sx = usable_w / max(layout["width"], 1)
    sy = usable_h / max(layout["height"] + 30, 1)
    scale = min(sx, sy)

    def tx(v: float) -> int:
        return margin + int(v * scale)

    def ty(v: float) -> int:
        return top_title + int(v * scale)

    def tw(v: float) -> int:
        return max(1, int(v * scale))

    def add_text(x, y, w, h, text, size=9.2, bold=False, align="center", opacity=None, color=None):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.clear()
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = _hex_rgb(color or theme["text"])
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        return tb

    def style_shape(sh, fill, stroke=None, line_width=1.1, transparency=0):
        sh.fill.solid()
        sh.fill.fore_color.rgb = _hex_rgb(fill)
        try:
            sh.fill.transparency = transparency
        except Exception:
            pass
        sh.line.color.rgb = _hex_rgb(stroke or theme["stroke"])
        sh.line.width = Pt(line_width)

    def label_node(node, x, y, w, h, ratio=0.67, size=9.0, subtitle=True):
        label = str(node.get("label", ""))
        add_text(x + int(w*0.05), y + int(h*ratio), int(w*0.90), int(h*0.23), label, size=size, bold=True)
        sub = node.get("subtitle") or node.get("shape")
        if subtitle and sub:
            add_text(x + int(w*0.05), y + int(h*0.86), int(w*0.90), int(h*0.13), str(sub), size=6.8, bold=False)

    def draw_scientific_illustration(ill, x, y, w, h):
        """Draw a small editable scientific illustration with native PPT shapes."""
        typ = str((ill or {}).get("type", ""))
        accent = theme.get("novel", "#D97706")
        blue = theme.get("input", "#2563EB")
        green = theme.get("output", "#4D7C0F")
        stroke = theme["stroke"]

        def circ(cx, cy, r, fill="#FFFFFF", outline=None, lw=0.7):
            sh=slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx-r), int(cy-r), int(2*r), int(2*r))
            style_shape(sh, fill, outline or stroke, lw)
            return sh

        def conn(x1,y1,x2,y2,lw=.6,dashed=False,color=None):
            c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2))
            c.line.color.rgb=_hex_rgb(color or stroke); c.line.width=Pt(lw)
            if dashed: c.line.dash_style=MSO_LINE_DASH_STYLE.DASH
            return c

        if typ == "custom_dsl":
            for p0 in ill.get("primitives") or []:
                shp=p0.get("shape")
                def X(v): return x+int(w*float(v)/100)
                def Y(v): return y+int(h*float(v)/100)
                fill0=p0.get("fill", "#FFFFFF")
                if fill0=="accent": fill0=accent
                elif fill0=="input": fill0=blue
                elif fill0=="output": fill0=green
                elif fill0 in {"none", None}: fill0="#FFFFFF"
                stroke0=p0.get("stroke", stroke)
                if stroke0=="accent": stroke0=accent
                elif stroke0=="input": stroke0=blue
                elif stroke0=="output": stroke0=green
                elif stroke0=="text": stroke0=theme["text"]
                if shp=="circle":
                    r=min(w,h)*float(p0.get("r",8))/100; circ(X(p0.get("cx",50)),Y(p0.get("cy",50)),r,fill0,stroke0,float(p0.get("stroke_width",.7)))
                elif shp=="ellipse":
                    rx=w*float(p0.get("rx",10))/100; ry=h*float(p0.get("ry",7))/100
                    sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,int(X(p0.get("cx",50))-rx),int(Y(p0.get("cy",50))-ry),int(2*rx),int(2*ry)); style_shape(sh,fill0,stroke0,float(p0.get("stroke_width",.7)))
                elif shp=="rect":
                    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,X(p0.get("x",10)),Y(p0.get("y",10)),int(w*float(p0.get("w",20))/100),int(h*float(p0.get("h",20))/100)); style_shape(sh,fill0,stroke0,float(p0.get("stroke_width",.7)))
                elif shp=="line":
                    conn(X(p0.get("x1",0)),Y(p0.get("y1",0)),X(p0.get("x2",100)),Y(p0.get("y2",100)),float(p0.get("stroke_width",.7)),bool(p0.get("dash")),stroke0)
                elif shp=="text":
                    add_text(X(p0.get("x",50))-int(w*.12),Y(p0.get("y",50))-int(h*.07),int(w*.24),int(h*.14),str(p0.get("text","")),size=float(p0.get("font_size",6.0)),color=stroke0)
            return

        if typ in {"timeseries", "sequence"}:
            conn(x+w*.12,y+h*.78,x+w*.90,y+h*.78,.55); conn(x+w*.12,y+h*.78,x+w*.12,y+h*.16,.55)
            pts=[]
            for i in range(18):
                px=x+w*(.14+.72*i/17); py=y+h*(.52-.14*math.sin(i*.83)-.05*math.sin(i*1.7)); pts.append((px,py))
            for a,b in zip(pts,pts[1:]): conn(*a,*b,.9,color=blue)
            return
        if typ == "feature_grid":
            for j in range(3,-1,-1):
                sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,int(x+w*(.16+.04*j)),int(y+h*(.18-.04*j)),int(w*.56),int(h*.55)); style_shape(sh,blue,stroke,.55,transparency=72-j*6)
            return
        if typ in {"globe_coordinates","fibonacci_sphere","spectral_sphere","field_map"}:
            cx=x+w*.50; cy=y+h*.47; r=min(w,h)*.30; circ(cx,cy,r,"#FFFFFF",stroke,.65)
            # editable latitude/longitude ellipses
            for fr in (.45,.72):
                sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,int(cx-r*fr),int(cy-r),int(2*r*fr),int(2*r)); style_shape(sh,"#FFFFFF",stroke,.35,transparency=100)
            for off in (-.45,0,.45):
                sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,int(cx-r*.78),int(cy+r*off-r*.08),int(2*r*.78),int(r*.16)); style_shape(sh,"#FFFFFF",stroke,.35,transparency=100)
            if typ=="fibonacci_sphere":
                n=int((ill.get("params") or {}).get("nodes",18)); phi=(1+5**.5)/2
                for i in range(n):
                    z=1-2*(i+.5)/n; rr=(1-z*z)**.5; th=2*math.pi*i/phi; px=cx+r*.82*rr*math.cos(th); py=cy-r*.82*z; circ(px,py,max(1.2,r*.025),accent,accent,.2)
            elif typ=="spectral_sphere":
                for k,col in enumerate((accent,green,blue)):
                    pts=[]
                    for i in range(16): pts.append((cx-r*.74+2*r*.74*i/15,cy+(k-1)*r*.25+r*.06*math.sin(i/15*math.pi*4+k)))
                    for a,b in zip(pts,pts[1:]): conn(*a,*b,.75,color=col)
            elif typ=="field_map":
                for off,col in ((-.25,blue),(0,green),(.25,accent)):
                    sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,int(cx-r*.72),int(cy+r*off-r*.10),int(2*r*.72),int(r*.20)); style_shape(sh,col,col,.1,transparency=75)
            elif typ=="globe_coordinates":
                circ(cx+r*.38,cy-r*.08,r*.07,accent,accent,.2); conn(cx+r*.38,cy,cx+r*.38,cy+r*.25,.65,color=accent)
            return
        if typ == "coordinate_axes":
            ox=x+w*.35; oy=y+h*.70; conn(ox,oy,x+w*.82,oy,.7); conn(ox,oy,ox,y+h*.18,.7); conn(ox,oy,x+w*.15,y+h*.84,.7); return
        if typ == "graph_network":
            pts=[(.18,.55),(.34,.28),(.50,.50),(.68,.22),(.78,.58),(.48,.78),(.28,.76)]; xy=[(x+w*a,y+h*b) for a,b in pts]
            for a,b in ((0,1),(1,2),(2,3),(3,4),(4,5),(5,6),(6,0),(1,6),(2,5),(2,4)): conn(*xy[a],*xy[b],.55,dashed=((a+b)%3==0))
            for i,(cx,cy) in enumerate(xy): circ(cx,cy,min(w,h)*.04,accent if i in {2,5} else "#FFFFFF",stroke,.6)
            return
        if typ == "attention_fan":
            q=(x+w*.50,y+h*.76); targets=[(x+w*.18,y+h*(.20+i*.16)) for i in range(3)]+[(x+w*.82,y+h*(.20+i*.16)) for i in range(3)]
            for i,pt in enumerate(targets): conn(*q,*pt,.5,True); circ(pt[0],pt[1],min(w,h)*.032,blue if i<3 else green,stroke,.5)
            circ(q[0],q[1],min(w,h)*.042,accent,stroke,.6); return
        if typ == "uncertainty_curve":
            conn(x+w*.12,y+h*.78,x+w*.90,y+h*.78,.5); pts=[]
            for i in range(22):
                t=-3+6*i/21; val=math.exp(-.5*t*t); pts.append((x+w*(.14+.72*i/21),y+h*(.75-.50*val)))
            for a,b in zip(pts,pts[1:]): conn(*a,*b,.9,color=blue)
            for fr in (.38,.62): conn(x+w*fr,y+h*.47,x+w*fr,y+h*.79,.45,True); return
        if typ == "satellite_occultation":
            sx=x+w*.24; sy=y+h*.24
            body=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,int(sx),int(sy),int(w*.13),int(h*.10)); style_shape(body,blue,stroke,.55,transparency=55)
            for d in (-1,1):
                px=sx+(d*w*.12)-(w*.08 if d<0 else 0); sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,int(px),int(sy-h*.01),int(w*.08),int(h*.12)); style_shape(sh,blue,stroke,.45,transparency=72)
            # Earth limb approximation + occultation rays
            earth=slide.shapes.add_shape(MSO_SHAPE.ARC,int(x+w*.10),int(y+h*.55),int(w*.82),int(h*.52)); earth.line.color.rgb=_hex_rgb(green); earth.line.width=Pt(1.2); earth.fill.background()
            for i in range(3): conn(sx+w*.06,sy+h*.08,x+w*(.69+i*.07),y+h*(.73+i*.045),.45,True)
            return
        if typ == "radar_dish":
            cx=x+w*.28; cy=y+h*.43; rr=min(w,h)*.18
            dish=slide.shapes.add_shape(MSO_SHAPE.ARC,int(cx-rr),int(cy-rr*.6),int(2*rr),int(1.35*rr)); dish.line.color.rgb=_hex_rgb(stroke); dish.line.width=Pt(.8); dish.fill.background()
            conn(cx,cy+rr*.4,cx-w*.03,y+h*.80,.75); conn(cx,cy+rr*.4,cx+w*.12,y+h*.80,.75)
            for j in range(3): conn(cx+rr*.45,cy-rr*.35,cx+rr*(1.15+j*.30),cy-rr*(.85-j*.05),.5,False,blue)
            return
        if typ == "sun_geomagnetic":
            scx=x+w*.25; scy=y+h*.45; sr=min(w,h)*.11; circ(scx,scy,sr,"#FFF4CC",accent,.7)
            for i in range(8):
                a=i*math.pi/4; conn(scx+math.cos(a)*sr*1.35,scy+math.sin(a)*sr*1.35,scx+math.cos(a)*sr*1.85,scy+math.sin(a)*sr*1.85,.55,False,accent)
            ecx=x+w*.68; ecy=y+h*.47; er=min(w,h)*.09; circ(ecx,ecy,er,"#FFFFFF",blue,.65)
            for sgn in (-1,1):
                for off in (.0,.14,.28): conn(ecx,ecy-er*(1.6+off),ecx+sgn*er*(1.7+off),ecy+er*(1.2+off),.35,False,blue)
            return
        if typ == "point_cloud":
            for i,(a,b) in enumerate(((.18,.28),(.30,.60),(.43,.36),(.55,.70),(.68,.25),(.80,.54),(.60,.45),(.35,.78),(.73,.75))): circ(x+w*a,y+h*b,min(w,h)*(.018+.004*(i%3)),accent if i%3==0 else blue,stroke,.3)
            return

    if fig.get("show_title", False) and fig.get("title"):
        add_text(margin, Inches(0.02), usable_w, Inches(0.28), str(fig["title"]), size=11.2, bold=True)

    # Build absolute bounding boxes first so edges can be placed behind visual glyphs.
    pos_abs: dict[str, tuple[int,int,int,int]] = {}
    panel_defs = []
    node_panel_direction = {}
    for pl in layout["panels"]:
        xoff = pl.get("x_offset", 0.0)
        yoff = pl.get("y_offset", 0.0) + 30
        panel_defs.append((pl, xoff, yoff))
        pdir = str((pl.get("panel") or {}).get("direction") or fig.get("direction", "LR"))
        for nid, box0 in pl["positions"].items():
            node_panel_direction[nid] = pdir
            box = {**box0, "x": box0["x"] + xoff, "y": box0["y"] + yoff}
            pos_abs[nid] = (tx(box["x"]), ty(box["y"]), tw(box["w"]), tw(box["h"]))

    # Preserve the obstacle-aware routes produced by the layout engine.  The route
    # remains a sequence of native PowerPoint connectors, not a raster overlay.
    route_abs: dict[tuple, list[tuple[int,int]]] = {}
    for pl, xoff, yoff in panel_defs:
        for e in pl.get("edges", []):
            pts=e.get("_route_points") or []
            if pts:
                key=(e.get("from"),e.get("to"),e.get("label",""),e.get("type","main"))
                route_abs[key]=[(tx(float(q[0])+xoff),ty(float(q[1])+yoff)) for q in pts]

    # Stage/group containers first so they remain behind connectors and node glyphs.
    for pl, xoff, yoff in panel_defs:
        for gb in pl.get("groups", []):
            gx, gy, gw, gh = tx(gb["x"]+xoff), ty(gb["y"]+yoff), tw(gb["w"]), tw(gb["h"])
            stage_boxes = bool(spec.get("style", {}).get("stage_containers")) and fig.get("layout_preset") == "publication_framework"
            if stage_boxes:
                role = str(gb.get("accent", "neutral"))
                fill_role = theme.get(role, theme["neutral"])
                heading = theme.get(role + "_text", theme["text"])
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, gw, gh)
                style_shape(sh, fill_role, heading, 0.75, transparency=90)
                add_text(gx+Inches(0.06), gy+Inches(0.015), gw-Inches(0.12), Inches(0.22), str(gb["label"]), size=7.2, bold=True, align="left", color=heading)
            else:
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, gw, gh)
                style_shape(sh, "#FAFBFC", theme["panel"], 0.7)
                sh.fill.transparency = 35
                sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                add_text(gx+Inches(0.05), gy+Inches(0.02), gw-Inches(0.1), Inches(0.14), str(gb["label"]).upper(), size=5.8, bold=True, align="left")

    # High-fan-in fusion nodes use one shared bus.  This preserves every input
    # relationship while avoiding the line-bundle clutter produced by drawing 3–6
    # independent arrowheads into the same small target.
    bundled_keys=set()
    incoming_by_target={}
    for e in spec.get("edges", []):
        incoming_by_target.setdefault(e.get("to"),[]).append(e)
    boxes_abs={nid:{"x":x,"y":y,"w":w,"h":h} for nid,(x,y,w,h) in pos_abs.items()}
    for tid,ies in incoming_by_target.items():
        node=node_map.get(tid,{})
        vtype=(node.get("visual") or {}).get("type")
        if len(ies)<3 or vtype not in {"fusion_node_publication","fusion_bar_publication"} or any(e.get("label") for e in ies):
            continue
        geo=fanin_bundle_geometry(ies,boxes_abs,tid,offset=tw(22.0))
        if not geo:
            continue
        for br in geo["branches"]:
            pts=br["points"]
            for (x1,y1),(x2,y2) in zip(pts,pts[1:]):
                c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2))
                c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(.95)
            e=br["edge"]; bundled_keys.add((e.get("from"),e.get("to"),e.get("label",""),e.get("type","main")))
        (x1,y1),(x2,y2)=geo["bus"]
        c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2)); c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(.95)
        pts=geo["trunk"]
        for j,((x1,y1),(x2,y2)) in enumerate(zip(pts,pts[1:])):
            c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2)); c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(1.35)
            if j==len(pts)-2:
                ln=c.line._get_or_add_ln(); tail=OxmlElement("a:tailEnd"); tail.set("type","triangle"); tail.set("w","sm"); tail.set("len","sm"); ln.append(tail)

    # High-fan-out encoders/attention blocks use one shared departure bus after any
    # fusion-specific fan-in edges have already been claimed.
    outgoing_by_source={}
    for e in spec.get("edges", []):
        outgoing_by_source.setdefault(e.get("from"),[]).append(e)
    for sid,oes in outgoing_by_source.items():
        remaining=[e for e in oes if (e.get("from"),e.get("to"),e.get("label",""),e.get("type","main")) not in bundled_keys]
        node=node_map.get(sid,{})
        vtype=(node.get("visual") or {}).get("type")
        if len(remaining)<3 or vtype not in {"encoder_module","emphasis_module","attention_block","feature_vector"} or any(e.get("label") for e in remaining):
            continue
        geo=fanout_bundle_geometry(remaining,boxes_abs,sid,offset=tw(22.0))
        if not geo:
            continue
        pts=geo["trunk"]
        for (x1,y1),(x2,y2) in zip(pts,pts[1:]):
            c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2)); c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(1.35)
        (x1,y1),(x2,y2)=geo["bus"]
        c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2)); c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(.95)
        for br in geo["branches"]:
            pts=br["points"]
            for j,((x1,y1),(x2,y2)) in enumerate(zip(pts,pts[1:])):
                c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,int(x1),int(y1),int(x2),int(y2)); c.line.color.rgb=_hex_rgb(theme["edge"]); c.line.width=Pt(.95)
                if j==len(pts)-2:
                    ln=c.line._get_or_add_ln(); tail=OxmlElement("a:tailEnd"); tail.set("type","triangle"); tail.set("w","sm"); tail.set("len","sm"); ln.append(tail)
            e=br["edge"]; bundled_keys.add((e.get("from"),e.get("to"),e.get("label",""),e.get("type","main")))

    # Connectors next: above stage backgrounds, below node glyphs.
    for e in spec.get("edges", []):
        if e.get("from") not in pos_abs or e.get("to") not in pos_abs:
            continue
        ax, ay, aw, ah = pos_abs[e["from"]]
        bx, by, bw, bh = pos_abs[e["to"]]
        direction = node_panel_direction.get(e.get("from"), fig.get("direction", "LR"))
        et = e.get("type", "main")
        key=(e.get("from"),e.get("to"),e.get("label",""),et)
        if key in bundled_keys:
            continue
        routed=route_abs.get(key)
        connectors=[]
        if routed and et != "residual":
            for j,((x1,y1),(x2,y2)) in enumerate(zip(routed,routed[1:])):
                conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2)
                connectors.append((conn,j==len(routed)-2,x1,y1,x2,y2))
        else:
            if et == "residual":
                x1, y1 = ax + int(aw*0.72), ay
                x2, y2 = bx + int(bw*0.28), by
                conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
            elif direction == "LR":
                x1, y1 = ax + aw, ay + ah // 2
                x2, y2 = bx, by + bh // 2
                conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
            else:
                x1, y1 = ax + aw // 2, ay + ah
                x2, y2 = bx + bw // 2, by
                conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
            connectors=[(conn,True,x1,y1,x2,y2)]
        for conn,is_last,x1,y1,x2,y2 in connectors:
            conn.line.color.rgb = _hex_rgb(theme["edge"])
            conn.line.width = Pt(1.35 if et == "main" else 0.95)
            if et in {"auxiliary", "conditioning", "training"}:
                conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            if is_last:
                ln = conn.line._get_or_add_ln()
                tail = OxmlElement("a:tailEnd")
                tail.set("type", "triangle")
                tail.set("w", "sm")
                tail.set("len", "sm")
                ln.append(tail)
        if e.get("label"):
            if routed and len(routed)>=2:
                segs=[]
                for p1,p2 in zip(routed,routed[1:]):
                    segs.append((abs(p2[0]-p1[0])+abs(p2[1]-p1[1]),p1,p2))
                _,p1,p2=max(segs,key=lambda z:z[0]); mx,my=(p1[0]+p2[0])//2,(p1[1]+p2[1])//2
            else:
                mx, my = (x1+x2)//2, (y1+y2)//2
            add_text(mx-Inches(0.33), my-Inches(0.13), Inches(0.66), Inches(0.18), e["label"], size=6.2)

    # Panel headers (if present) sit above connectors.
    for pl, xoff, yoff in panel_defs:
        pdef = pl["panel"]
        if pdef.get("label") or pdef.get("title"):
            add_text(tx(xoff+8), ty(yoff+4), tw(pl["width"]-16), Inches(0.22), f"{pdef.get('label','')} {pdef.get('title','')}".strip(), size=9.2, bold=True, align="left")

    # Scientific visual glyphs. All primitives remain independently editable.
    for nid, (x, y, w, h) in pos_abs.items():
        node = node_map[nid]
        role = node.get("role", "neutral")
        fill = theme.get(role, theme["neutral"])
        stroke = theme["stroke"]
        visual = node.get("visual") or {}
        vtype = visual.get("type", "generic_module") if isinstance(visual, dict) else str(visual)
        ill = node.get("illustration") or None

        if ill and vtype != "input_card_publication":
            card=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
            style_shape(card,"#FFFFFF",stroke,1.05 if node.get("role")=="novel" else .9)
            band=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,int(h*.055)); style_shape(band,fill,fill,.2)
            comp=str(ill.get("composition","illustration-top"))
            if comp=="illustration-left":
                ix,iy,iw,ih=x+int(w*.05),y+int(h*.18),int(w*.31),int(h*.62); draw_scientific_illustration(ill,ix,iy,iw,ih)
                add_text(x+int(w*.43),y+int(h*.24),int(w*.52),int(h*.34),node.get("label",""),size=8.2,bold=True,align="left")
                sub=node.get("subtitle") or node.get("shape");
                if sub: add_text(x+int(w*.43),y+int(h*.62),int(w*.52),int(h*.18),sub,size=6.2,align="left")
            elif comp=="illustration-bottom":
                draw_scientific_illustration(ill,x+int(w*.18),y+int(h*.57),int(w*.64),int(h*.34))
                add_text(x+int(w*.05),y+int(h*.16),int(w*.90),int(h*.30),node.get("label",""),size=8.2,bold=True)
            else:
                draw_scientific_illustration(ill,x+int(w*.16),y+int(h*.10),int(w*.68),int(h*.43))
                add_text(x+int(w*.05),y+int(h*.68),int(w*.90),int(h*.20),node.get("label",""),size=8.2,bold=True)
                sub=node.get("subtitle") or node.get("shape");
                if sub: add_text(x+int(w*.08),y+int(h*.86),int(w*.84),int(h*.10),sub,size=6.0)

        elif vtype == "input_card_publication":
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            style_shape(card, "#FFFFFF", stroke, 0.95)
            symbol = visual.get("input_symbol", "vector")
            vx, vy = x+int(w*0.08), y+int(h*0.25)
            vw, vh = int(w*0.24), int(h*0.40)
            if ill:
                draw_scientific_illustration(ill, vx-int(w*.02), y+int(h*.16), int(w*.29), int(h*.64))
            elif symbol == "tensor":
                for i in range(2,-1,-1):
                    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, vx+int(i*w*0.018), vy-int(i*h*0.025), vw, vh)
                    style_shape(sh, fill, stroke, 0.65, transparency=58-i*8)
            elif symbol == "tokens":
                gap=max(1,int(w*0.008)); bw=max(2,int((vw-gap*4)/5))
                for i in range(5):
                    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vx+i*(bw+gap), vy+int(vh*0.12), bw, int(vh*0.76))
                    style_shape(sh, fill, stroke, 0.55, transparency=45)
            else:
                heights=[.35,.70,.50,.84,.42]
                bw=max(2,int(vw*.12)); gap=max(1,int(vw*.07))
                for i,fr in enumerate(heights):
                    bh=int(vh*fr)
                    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, vx+i*(bw+gap), vy+vh-bh, bw, bh)
                    style_shape(sh, fill, stroke, 0.45, transparency=48)
            tx0=x+int(w*0.40); tw0=int(w*0.55)
            add_text(tx0, y+int(h*0.20), tw0, int(h*0.40), node.get("label", ""), size=8.4, bold=True, align="left")
            if node.get("shape"):
                add_text(tx0, y+int(h*0.62), tw0, int(h*0.20), node.get("shape"), size=6.4, align="left")

        elif vtype in {"data_stack", "feature_map_stack", "modality_card", "unet_stage", "bottleneck", "input_tensor", "feature_tensor", "field_tensor", "field_tensor_output"}:
            layers = 4 if vtype in {"feature_map_stack", "bottleneck"} else 3
            cw, ch = int(w*0.67), int(h*0.48)
            bx, by = x + int(w*0.14), y + int(h*0.12)
            for i in range(layers-1, -1, -1):
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx+int(i*w*0.035), by-int(i*h*0.035), cw, ch)
                style_shape(sh, fill, stroke, 0.8, transparency=min(55, 18+i*9))
            if vtype == "unet_stage":
                # Editable scale-direction triangle.
                is_dec = "decoder" in str(node.get("label", "")).lower()
                shp = MSO_SHAPE.UP_ARROW if is_dec else MSO_SHAPE.DOWN_ARROW
                ar = slide.shapes.add_shape(shp, x+int(w*0.79), y+int(h*0.18), int(w*0.08), int(h*0.18))
                style_shape(ar, stroke, stroke, 0.4, transparency=35)
            label_node(node, x, y, w, h, ratio=0.70, size=8.6)

        elif vtype in {"token_strip", "sequence_strip", "token_matrix", "token_strip_publication"}:
            outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+int(w*0.08), y+int(h*0.10), int(w*0.84), int(h*0.44))
            style_shape(outer, fill, stroke, 0.9)
            rows = 3 if vtype == "token_matrix" else 1
            cols = 7
            ix, iy, iw, ih = x+int(w*0.12), y+int(h*0.16), int(w*0.76), int(h*0.32)
            gap = int(w*0.012)
            cw = max(2, int((iw-gap*(cols-1))/cols))
            ch = max(2, int((ih-gap*(rows-1))/rows))
            for rr in range(rows):
                for cc in range(cols):
                    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ix+cc*(cw+gap), iy+rr*(ch+gap), cw, ch)
                    style_shape(sh, stroke, stroke, 0.2, transparency=82)
            label_node(node, x, y, w, h, ratio=0.66, size=8.5)

        elif vtype == "embedding_card":
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            style_shape(sh, fill, stroke, 1.0)
            for i in range(5):
                bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+int(w*0.10), y+int(h*(0.12+0.07*i)), int(w*(0.50+0.07*(i%3))), int(h*0.035))
                style_shape(bar, stroke, stroke, 0.2, transparency=82)
            label_node(node, x, y, w, h, ratio=0.63, size=8.4)

        elif vtype in {"transformer_stack", "transformer_macro"}:
            bw, bh = int(w*0.72), int(h*0.47)
            bx, by = x+int(w*0.12), y+int(h*0.13)
            layers = min(4, max(3, int(visual.get("layers", 3))))
            for i in range(layers-1, -1, -1):
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx+int(i*w*0.035), by-int(i*h*0.03), bw, bh)
                style_shape(sh, fill, stroke, 0.8, transparency=min(52, 16+i*9))
            label_node(node, x, y, w, h, ratio=0.68, size=8.7)

        elif vtype in {"attention_heads", "attention_block"}:
            heads = min(6, max(3, int(visual.get("heads", 4))))
            target_x, target_y = x+w//2, y+int(h*0.50)
            for i in range(heads):
                hx = x+int(w*(0.18+0.64*i/(heads-1)))
                hy = y+int(h*0.19)
                circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, hx-int(w*0.045), hy-int(w*0.045), int(w*0.09), int(w*0.09))
                style_shape(circ, fill, stroke, 0.8)
                conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, hx, hy+int(w*0.045), target_x, target_y)
                conn.line.color.rgb = _hex_rgb(stroke); conn.line.width = Pt(0.65)
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, target_x-int(w*0.17), target_y-int(h*0.04), int(w*0.34), int(h*0.08))
            style_shape(bar, stroke, stroke, 0.3, transparency=80)
            label_node(node, x, y, w, h, ratio=0.67, size=8.2)

        elif vtype in {"ffn_block", "ffn_block_publication"}:
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            style_shape(sh, fill, stroke, 1.0)
            for row in range(2):
                widths = [0.17, 0.25, 0.17]
                xx = x+int(w*0.16)
                for j, frac in enumerate(widths):
                    bw = int(w*frac)
                    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xx, y+int(h*(0.14+0.16*row)), bw, int(h*0.07))
                    style_shape(bar, stroke, stroke, 0.2, transparency=84-j*4)
                    xx += bw+int(w*0.025)
            label_node(node, x, y, w, h, ratio=0.64, size=8.4)

        elif vtype in {"norm_bar", "norm_bar_publication"}:
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+int(w*0.05), y+int(h*0.14), int(w*0.90), int(h*0.25))
            style_shape(sh, fill, stroke, 0.95)
            add_text(x, y+int(h*0.54), w, int(h*0.28), node.get("label", "Norm"), size=6.5, bold=True)

        elif vtype in {"graph_input", "graph_message"}:
            outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            style_shape(outer, fill, stroke, 0.9, transparency=7)
            pts = [(0.25,0.35),(0.38,0.20),(0.55,0.25),(0.72,0.39),(0.54,0.52),(0.33,0.53)]
            coords=[(x+int(w*a),y+int(h*b)) for a,b in pts]
            for a,b in [(0,1),(1,2),(2,3),(3,4),(4,5),(5,0),(1,5),(2,4)]:
                conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,*coords[a],*coords[b])
                conn.line.color.rgb=_hex_rgb(stroke); conn.line.width=Pt(0.6)
            for i,(cx,cy) in enumerate(coords):
                rr=int(w*(0.032 if vtype=="graph_input" else (0.042 if i in {2,4} else 0.032)))
                circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,cx-rr,cy-rr,rr*2,rr*2)
                style_shape(circ,"#FFFFFF",stroke,0.9)
            label_node(node,x,y,w,h,ratio=0.68,size=8.1)

        elif vtype == "graph_pool":
            center=(x+w//2,y+int(h*0.50))
            for i in range(5):
                cx=x+int(w*(0.25+0.125*i)); cy=y+int(h*0.22)
                circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,cx-int(w*0.025),cy-int(w*0.025),int(w*0.05),int(w*0.05))
                style_shape(circ,fill,stroke,0.7)
                conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,cx,cy,center[0],center[1])
                conn.line.color.rgb=_hex_rgb(stroke); conn.line.width=Pt(0.5)
            label_node(node,x,y,w,h,ratio=0.63,size=7.8)

        elif vtype == "router_gate":
            sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x+int(w*0.25), y+int(h*0.08), int(w*0.50), int(h*0.50))
            style_shape(sh, fill, stroke, 1.0)
            label_node(node,x,y,w,h,ratio=0.65,size=8.1)

        elif vtype == "expert_fan":
            count=min(6,max(3,int(visual.get("experts",4))))
            for i in range(count):
                sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x+int(w*(0.18+0.02*(i%2))),y+int(h*(0.07+0.12*i)),int(w*0.56),int(h*0.105))
                style_shape(sh,fill,stroke,0.7,transparency=max(5,30-i*3))
                add_text(x+int(w*0.22),y+int(h*(0.075+0.12*i)),int(w*0.12),int(h*0.08),f"E{i+1}",size=5.8)
            total=int(visual.get("total_experts",node.get("repeat",1) or 1))
            if total>count:
                add_text(x+int(w*0.76),y+int(h*0.37),int(w*0.20),int(h*0.15),f"… ×{total}",size=7.2,bold=True)
            label_node(node,x,y,w,h,ratio=0.84,size=7.8,subtitle=False)

        elif vtype == "fusion_bar_publication":
            sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
            style_shape(sh,"#FFFFFF",stroke,0.95)
            band=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,int(w*0.035),h)
            style_shape(band,theme.get("fusion",fill),theme.get("fusion",fill),0.2)
            add_text(x+int(w*0.08),y+int(h*0.14),int(w*0.86),int(h*0.44),node.get("label","Feature fusion"),size=7.8,bold=True)
            sub=node.get("subtitle") or "Concatenate"
            add_text(x+int(w*0.08),y+int(h*0.60),int(w*0.86),int(h*0.22),sub,size=6.2)

        elif vtype in {"merge_glyph", "weighted_merge", "add_node", "fusion_node_publication"}:
            shape_type = MSO_SHAPE.DIAMOND if vtype=="weighted_merge" else MSO_SHAPE.OVAL
            sh=slide.shapes.add_shape(shape_type,x+int(w*0.18),y+int(h*0.12),int(w*0.64),int(h*0.64))
            style_shape(sh,fill,stroke,1.1)
            symbol="Sw" if vtype=="weighted_merge" else ("C" if vtype=="fusion_node_publication" else str(visual.get("symbol",node.get("label","+"))))
            add_text(x+int(w*0.18),y+int(h*0.12),int(w*0.64),int(h*0.64),symbol,size=9.0,bold=True)

        elif vtype == "recurrent_cells":
            cells=min(5,max(3,int(visual.get("cells",4))))
            gap=int(w*0.025); cw=int((w*0.72-gap*(cells-1))/cells); sx=x+int(w*0.14); cy=y+int(h*0.18)
            for i in range(cells):
                sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,sx+i*(cw+gap),cy,cw,int(h*0.27))
                style_shape(sh,fill,stroke,0.8)
            label_node(node,x,y,w,h,ratio=0.65,size=8.2)

        elif vtype in {"spectral_operator", "spectral_operator_publication"}:
            sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); style_shape(sh,fill,stroke,1.0)
            # Editable spectrum bars.
            heights=[0.08,0.16,0.26,0.15,0.08]
            for i,fr in enumerate(heights):
                bh=int(h*fr); bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x+int(w*(0.31+0.075*i)),y+int(h*0.48)-bh,int(w*0.035),bh)
                style_shape(bar,stroke,stroke,0.2,transparency=78)
            label_node(node,x,y,w,h,ratio=0.66,size=8.1)

        elif vtype == "fusion_hub":
            cx,cy=x+w//2,y+int(h*0.40)
            for ang in (-150,-90,-30):
                a=math.radians(ang); px=cx+int(w*0.25*math.cos(a)); py=cy+int(h*0.25*math.sin(a))
                dot=slide.shapes.add_shape(MSO_SHAPE.OVAL,px-int(w*0.035),py-int(w*0.035),int(w*0.07),int(w*0.07)); style_shape(dot,fill,stroke,0.7)
                conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,px,py,cx,cy); conn.line.color.rgb=_hex_rgb(stroke); conn.line.width=Pt(0.6)
            hub=slide.shapes.add_shape(MSO_SHAPE.OVAL,cx-int(w*0.09),cy-int(w*0.09),int(w*0.18),int(w*0.18)); style_shape(hub,fill,stroke,1.0)
            label_node(node,x,y,w,h,ratio=0.64,size=7.7)

        else:
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
            if vtype == "output_card": shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
            sh = slide.shapes.add_shape(shape_type, x, y, w, h)
            style_shape(sh, fill, stroke, 1.15 if vtype=="output_card" else 1.0)
            label_node(node,x,y,w,h,ratio=0.38 if vtype=="output_card" else 0.35,size=8.7)

        repeat=int(node.get("repeat",1) or 1)
        if repeat>1 and vtype != "expert_fan":
            bw,bh=int(w*0.20),int(h*0.16)
            badge=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x+w-bw-int(w*0.03),y+int(h*0.04),bw,bh)
            style_shape(badge,"#FFFFFF",stroke,0.6)
            add_text(badge.left,badge.top,badge.width,badge.height,f"×{repeat}",size=6.5,bold=True)
        ref=node.get("detail_panel_ref")
        if ref:
            add_text(x+int(w*0.03), y+int(h*0.02), int(w*0.16), int(h*0.14), str(ref), size=6.5, bold=True, align="left", color=theme.get("novel_text", theme["text"]))

    prs.save(output)
    return output
