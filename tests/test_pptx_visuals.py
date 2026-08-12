from pathlib import Path
import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation

from ml_architecture_diagram.spec import load_spec
from ml_architecture_diagram.renderers.pptx import render_pptx

ROOT = Path(__file__).resolve().parents[1]


def test_pptx_uses_multiple_editable_shapes_for_visual_grammar(tmp_path):
    spec = load_spec(ROOT / "examples" / "transformer_encoder.yaml")
    out = render_pptx(spec, tmp_path / "transformer.pptx")
    prs = Presentation(out)
    assert len(prs.slides) == 1
    # More shapes than architecture nodes indicates composite native glyphs, not one raster per node.
    assert len(prs.slides[0].shapes) > len(spec["nodes"]) * 2
    assert all(not getattr(shape, "shape_type", None) == 13 for shape in prs.slides[0].shapes)  # 13 = picture in python-pptx enum
